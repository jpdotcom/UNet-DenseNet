from pptx import Presentation
from pptx.util import Inches

from datetime import date
from skimage import color
import glob
import cv2
import numpy as np
from io import BytesIO
import PIL.Image as Image
from Test import getImages

title = "Tubole segmentation"
today = date.today()
author = "Jay Patel"
comments = "data and code taken from blog andrewjanowczyk.com "
pptxfname = "TuboleSegmentation.pptx"


prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(10)
 
blank_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(blank_slide_layout)


slide.placeholders[0].text = title
 
tf = slide.placeholders[1].text_frame
tf.text = f'Date: {today}\n'
tf.text += f"Author: {author}\n"
tf.text += f"Comments: {comments}\n"

def blend2Images(img, mask):
    if (img.ndim == 3):
        img = color.rgb2gray(img)
    if (mask.ndim == 3):
        mask = color.rgb2gray(mask)
    img = img[:, :, None] * 1.0  # can't use boolean
    mask = mask[:, :, None] * 1.0
    out = np.concatenate((mask, img, mask), 2) * 255
    return out.astype('uint8')

def addimagetoslide(slide,img,left,top, height, width, resize = .1):
    res = cv2.resize(img , None, fx=resize,fy=resize ,interpolation=cv2.INTER_CUBIC) #since the images are going to be small, we can resize them to prevent the final pptx file from being large for no reason
    image_stream = BytesIO()
    Image.fromarray(res).save(image_stream,format="PNG")
 
    pic = slide.shapes.add_picture(image_stream, left, top ,height,width)
    image_stream.close()

def normalize(a):
    norm_image = cv2.normalize(a, None, alpha = 0, beta = 255, norm_type = cv2.NORM_MINMAX, dtype = cv2.CV_32F)

    norm_image = norm_image.astype(np.uint8)

    return norm_image

def blend2Images(img, mask):
    if (img.ndim == 3):
        img = color.rgb2gray(img)
    if (mask.ndim == 3):
        mask = color.rgb2gray(mask)
    img = img[:, :, None] * 1.0  # can't use boolean
    mask = mask[:, :, None] * 1.0
    out = np.concatenate((mask, img, mask), 2) * 255
    return out.astype('uint8')
if (__name__=="__main__"):

    total_img,total_labels,total_out=getImages();

    for i in range(len(total_img)):
        blank_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(blank_slide_layout)
        img,label,out=normalize(total_img[i]),normalize(total_labels[i]),normalize(total_out[i]);

        diff=blend2Images(out,label);
            

        addimagetoslide(slide, img, Inches(0),Inches(0),Inches(5),Inches(5))
        addimagetoslide(slide, label, Inches(5),Inches(5),Inches(5),Inches(5))
        addimagetoslide(slide, out, Inches(5),Inches(0),Inches(5),Inches(5))
        addimagetoslide(slide,diff, Inches(0),Inches(5),Inches(5),Inches(5))

prs.save(pptxfname)
